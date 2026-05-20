#!/usr/bin/env python3
"""Generate the 'Consumerism & Fast Fashion' PowerPoint deck and a Word
speaker-notes manuscript from a single content source.

Outputs (next to this script):
  - Fast_Fashion_Presentation.pptx
  - Fast_Fashion_Speaker_Notes.docx
"""

import math
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

from docx import Document
from docx.shared import Pt as DPt, RGBColor as DRGB, Inches as DInches
from docx.enum.text import WD_ALIGN_PARAGRAPH

HERE = os.path.dirname(os.path.abspath(__file__))

# ---------------------------------------------------------------------------
# Palette & fonts
# ---------------------------------------------------------------------------
PRIMARY = RGBColor(0x1E, 0x3D, 0x32)   # deep forest green
SECONDARY = RGBColor(0x3F, 0x7D, 0x6A)  # mid green
ACCENT = RGBColor(0xC0, 0x6A, 0x3C)    # clay / terracotta
BG = RGBColor(0xF6, 0xF3, 0xEC)        # warm off-white
INK = RGBColor(0x22, 0x2A, 0x26)       # near-black
MUTED = RGBColor(0x6B, 0x73, 0x6E)     # grey-green
MUTED_DARK = RGBColor(0x4A, 0x52, 0x4D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HEAD = "Georgia"
BODY = "Calibri"

# ---------------------------------------------------------------------------
# Content: every slide = tag (eyebrow), title, bullets, manuscript (notes)
# bullets: list of items; item is str (level 0) or (level, str)
# ---------------------------------------------------------------------------
SLIDES = [
    {
        "type": "title",
        "title": "Consumerism & Fast Fashion",
        "subtitle": "The hidden cost of cheap clothes",
        "meta": [
            "English – Group Presentation (makkerskab)",
            "By: [Your name] & [Your name]",
            "Deadline: 24/05-2026",
        ],
        "notes": (
            "Hello, and welcome to our presentation on consumerism and fast fashion. "
            "We are [Name] and [Name]. Today we will look at the hidden cost behind the "
            "cheap clothes that many of us buy. We base our discussion on two texts: the "
            "Earth.org article 'Fast Fashion: The Danger of Sweatshops' and the chapter "
            "'The Greening of Throwaway Stuff' from the textbook BusinessLike. Let's start "
            "with an overview of what we will cover."
        ),
    },
    {
        "tag": "Overview",
        "title": "Agenda",
        "bullets": [
            "Key terms: consumerism, materialism & fast fashion",
            "Our two texts",
            "Rhetorical analysis of “The Danger of Sweatshops” – the pentagram",
            "Modes of persuasion & other rhetorical devices",
            "Message, intention & the sender–receiver relationship",
            "Discussion: pros & cons of fast fashion",
            "Consumer responsibility – what can we do?",
            "Conclusion",
        ],
        "notes": (
            "Here is our agenda. First we will define three key terms – consumerism, "
            "materialism and fast fashion. Then we briefly introduce our two texts, before "
            "we move on to the main part: a rhetorical analysis of the Earth.org article "
            "using the rhetorical pentagram. After that we look at the modes of persuasion "
            "and other rhetorical devices, and comment on the message and intention of the "
            "article. Finally we discuss the pros and cons of fast fashion, the question of "
            "consumer responsibility, and what we can actually do – before we round off "
            "with a conclusion."
        ),
    },
    {
        "tag": "Background",
        "title": "Our two texts",
        "bullets": [
            "“Fast Fashion: The Danger of Sweatshops” – Earth.org",
            (1, "An environmental, non-profit news & data organisation"),
            (1, "Focus: the human & environmental cost of cheap clothing"),
            "“The Greening of Throwaway Stuff” – BusinessLike (Systime)",
            (1, "Focus: our throwaway culture & how it can become ‘greener’"),
            "Both deal with overconsumption – a case + a wider perspective",
        ],
        "notes": (
            "Our two texts approach the same problem from different angles. 'Fast Fashion: "
            "The Danger of Sweatshops' is published by Earth.org, an environmental "
            "organisation and news site, and it focuses on the human and environmental cost "
            "of cheap clothing. 'The Greening of Throwaway Stuff' from BusinessLike looks "
            "more broadly at our throwaway culture and how consumer society might become "
            "greener and more sustainable. Together they give us both a concrete case and a "
            "wider perspective on overconsumption."
        ),
    },
    {
        "tag": "Definitions",
        "title": "Key terms",
        "bullets": [
            "Consumerism – the society",
            (1, "A social & economic order that pushes us to keep buying goods"),
            "Materialism – the mindset",
            (1, "Owning things & money matter more than other values"),
            "Fast fashion – the product",
            (1, "Cheap, mass-produced, trend-copying clothes made to be thrown away"),
            (1, "e.g. Shein, H&M, Zara, Primark"),
        ],
        "notes": (
            "Before we analyse, let's define three terms. Consumerism is a social and "
            "economic order that encourages people to buy more and more goods – a society "
            "where consumption is seen as a path to happiness and status. Materialism is the "
            "individual mindset behind it: the belief that owning things and money matter "
            "more than, for example, relationships or experiences. Fast fashion is a concrete "
            "result of both: cheap, mass-produced clothing that copies catwalk trends at high "
            "speed and is designed to be bought and thrown away quickly – think Shein, "
            "H&M, Zara and Primark. So there is a chain here: consumerism is the society, "
            "materialism is the mindset, and fast fashion is the product."
        ),
    },
    {
        "type": "pentagram",
        "tag": "Rhetorical analysis",
        "title": "The rhetorical pentagram",
        "bullets": [
            "A model of the communication situation",
            "Five points that all influence each other",
            "Change one – and the others change too",
        ],
        "notes": (
            "To analyse the Earth.org article we use the rhetorical pentagram. The pentagram "
            "is a model of the communication situation, made up of five points that all "
            "influence each other: the sender, the receiver, the topic or subject, the "
            "language, and the circumstances or context. The key idea is that these five are "
            "connected – if you change one, the others change too. For example, who the "
            "sender is affects what language they use and how the receiver reacts. We will "
            "now go through all five points one by one."
        ),
    },
    {
        "tag": "Pentagram – 1 of 5",
        "title": "Sender (Afsender)",
        "bullets": [
            "Earth.org – an environmental, non-profit news & data organisation",
            "An institution, not a single famous author → ‘expert’ ethos",
            "Mission-driven: advocacy & education → not neutral, has an agenda",
            "Credibility built on data, surveys & a serious tone",
        ],
        "notes": (
            "First, the sender. The article comes from Earth.org, which is an environmental, "
            "non-profit news and data organisation. That is an important point: the sender is "
            "an institution, not a single famous journalist. This gives the article a kind of "
            "'expert organisation' ethos – it appears authoritative and well-researched. "
            "But we should also notice that Earth.org is mission-driven: its purpose is "
            "environmental advocacy and education. So the sender is not neutral – it has a "
            "clear agenda, and it builds its credibility through data, surveys and a serious, "
            "factual tone. If a specific author is shown on the page, mention their name here."
        ),
    },
    {
        "tag": "Pentagram – 2 of 5",
        "title": "Receiver (Modtager)",
        "bullets": [
            "Environmentally conscious, educated readers; students; the young",
            "Western/global consumers who actually buy fast fashion",
            "People open to changing their habits – or at least curious",
            "Addressed as both rational AND moral beings",
        ],
        "notes": (
            "Second, the receiver. The article is aimed primarily at environmentally "
            "conscious, educated readers – students, younger people, and the kind of "
            "consumers who actually buy fast fashion in the Western world. These are people "
            "who are open to changing their habits, or at least curious about the issue. The "
            "article addresses them as both rational beings – through facts and statistics "
            "– and as moral beings – through its focus on suffering and injustice. There "
            "is a slight risk of 'preaching to the choir' here, since the people most likely "
            "to read an Earth.org article may already agree."
        ),
    },
    {
        "tag": "Pentagram – 3 of 5",
        "title": "Topic / Subject (Emne)",
        "bullets": [
            "The human & environmental cost of fast-fashion sweatshops",
            "Exploitation: poverty wages, child labour, unsafe factories",
            (1, "Rana Plaza, Bangladesh, 2013 – 1,000+ workers died"),
            "Waste: 80–100 billion garments/year; a truckload burned every second",
            "The ‘true price’ hidden behind the cheap price tag",
        ],
        "notes": (
            "Third, the topic or subject. The article is about the human and environmental "
            "cost of fast-fashion sweatshops. On the human side, it describes exploitation: "
            "poverty wages, child labour, and unsafe factories – the clearest example being "
            "the Rana Plaza collapse in Bangladesh in 2013, where over a thousand garment "
            "workers died. On the environmental side, it points to enormous waste: between 80 "
            "and 100 billion new garments are made every year, and a truckload of clothing is "
            "burned or sent to landfill every single second. The core idea is that there is a "
            "'true price' hidden behind the cheap price tag."
        ),
    },
    {
        "tag": "Pentagram – 4 of 5",
        "title": "Language (Sprog)",
        "bullets": [
            "Formal but accessible English – academic yet readable",
            "Emotive, loaded words: ‘grueling’, ‘cruel’, ‘inhumane’, ‘exploited’",
            "Heavy use of facts, figures & surveys → objective register",
            "A deliberate blend of the factual and the emotional",
        ],
        "notes": (
            "Fourth, the language. The article is written in formal but accessible English "
            "– academic, but easy enough for a general reader. The word choices are often "
            "emotive and loaded: words like 'grueling', 'cruel', 'inhumane' and 'exploited'. "
            "At the same time, the text leans heavily on facts, figures and surveys, which "
            "gives it an objective register. So the language is a deliberate blend of the "
            "factual and the emotional – and that combination is exactly what makes it "
            "persuasive."
        ),
    },
    {
        "tag": "Pentagram – 5 of 5",
        "title": "Circumstances / Context (Omstændigheder)",
        "bullets": [
            "Written amid the global climate & sustainability debate",
            "The boom of ultra-fast fashion (e.g. Shein)",
            "Post-Rana Plaza awareness of garment-worker safety",
            "Published online – free, shareable, global, hyperlinked sources",
        ],
        "notes": (
            "Fifth, the circumstances or context. The article is written in the middle of a "
            "global climate and sustainability debate, and at a time when ultra-fast-fashion "
            "brands like Shein are booming. It also builds on a public awareness of "
            "garment-worker safety that grew after disasters like Rana Plaza. The fact that "
            "it is published online matters too: it is free, easily shareable, reaches a "
            "global audience, and can link directly to its sources. It is part of Earth.org's "
            "wider campaign of environmental journalism. As you can see, all five points "
            "connect – the green sender, the concerned receiver, the alarming topic, the "
            "emotive-but-factual language, and the climate-crisis context all reinforce one "
            "another."
        ),
    },
    {
        "tag": "Rhetoric",
        "title": "Modes of persuasion",
        "bullets": [
            "Ethos – Earth.org’s authority; sober, well-sourced tone",
            "Logos – statistics & cause–effect",
            (1, "60% of surveyed Indian mill workers were under 18"),
            (1, "wages cover only 1/5–1/2 of a family’s needs; 1,000+ Rana Plaza deaths"),
            "Pathos – emotive language; women & children; cruel conditions",
        ],
        "notes": (
            "Now to the modes of persuasion – ethos, logos and pathos. The article uses "
            "ethos through the authority of Earth.org as an environmental organisation, and "
            "through its sober, well-sourced tone. It uses logos through statistics and "
            "cause-and-effect reasoning – for example, that sixty percent of workers in "
            "surveyed Indian mills were under eighteen, that minimum wages cover only a fifth "
            "to a half of what a family needs, and the more than a thousand deaths at Rana "
            "Plaza. And it uses pathos through emotive language and its focus on women and "
            "children suffering in cruel conditions. The strongest effect comes from "
            "combining logos and pathos: the numbers make us think, and the human stories "
            "make us feel."
        ),
    },
    {
        "tag": "Rhetoric",
        "title": "Other rhetorical devices",
        "bullets": [
            "Loaded / connotative language → guilt & moral appeal",
            "Vivid imagery & near-hyperbole (“a truckload burned every second”)",
            "Contrast: our cheap clothes ↔ their high human cost",
            "Authority & evidence – surveys, named events",
        ],
        "notes": (
            "Beyond the three appeals, the article uses several other rhetorical devices. "
            "There is loaded, connotative language that creates a sense of guilt and moral "
            "responsibility. There is vivid imagery and almost hyperbolic framing – like 'a "
            "truckload of clothes burned every second' – which makes an abstract problem "
            "feel concrete and urgent. There is contrast or juxtaposition between our cheap, "
            "disposable clothes and the high human cost behind them. And there is a "
            "consistent use of authority and evidence – surveys and named events – to back "
            "up the emotional appeal."
        ),
    },
    {
        "tag": "Interpretation",
        "title": "Message & intention",
        "bullets": [
            "Message: the low price is paid by exploited workers & a damaged planet",
            (1, "the real cost is hidden from the consumer"),
            "Intention: to inform AND to persuade",
            (1, "raise awareness → a call to rethink consumption"),
            "Educates with logos + moves us with pathos",
        ],
        "notes": (
            "So what is the message and intention? The message is that the low price of fast "
            "fashion is actually paid by exploited workers and a damaged planet – the real "
            "cost is hidden from the consumer. The intention is twofold: to inform and to "
            "persuade. Earth.org wants to raise awareness, but it also wants to change "
            "behaviour – it is essentially a call to action to rethink how we consume. So the "
            "article educates us with logos and moves us with pathos at the same time."
        ),
    },
    {
        "tag": "Interpretation",
        "title": "Sender ↔ receiver: does the language match?",
        "bullets": [
            "Largely yes – an eco-organisation writing for value-driven readers",
            "Formal-but-accessible register fits educated, non-expert consumers",
            "Hard data + emotional appeal fits rational AND moral readers",
            "Small mismatch: very emotive words may feel one-sided to sceptics",
        ],
        "notes": (
            "The assignment asks us to consider the language between the sender and the "
            "receiver – does it match? For the most part, yes. An environmental organisation "
            "writing for informed, value-driven readers uses exactly the register we would "
            "expect: formal but accessible, mixing hard data with emotional appeal. That fits "
            "an audience that is both rational and morally engaged. The one small mismatch is "
            "that the very emotive vocabulary could feel one-sided to a more sceptical reader "
            "– but since the likely audience already cares about the environment, the "
            "language and the receiver are well matched overall."
        ),
    },
    {
        "tag": "Discussion",
        "title": "Pros of fast fashion",
        "bullets": [
            "Affordable & accessible – style for low-income consumers",
            "Provides jobs & income in developing countries",
            "Fast, varied, on-trend – drives economic activity",
            "Be fair: acknowledge the benefits before we judge",
        ],
        "notes": (
            "Now let's discuss fast fashion more broadly, starting with the pros – because it "
            "is not only negative. Fast fashion is affordable and accessible: it lets people "
            "on low incomes follow trends and express themselves through clothes. It also "
            "creates jobs and income in developing countries, where those jobs may be among "
            "the few available. And it is fast, varied and responsive to trends, which drives "
            "a lot of economic activity. We have to be fair and acknowledge these benefits "
            "before we judge."
        ),
    },
    {
        "tag": "Discussion",
        "title": "Cons of fast fashion",
        "bullets": [
            "Labour exploitation: poverty wages, child labour, unsafe factories",
            "Environmental harm: water pollution, textile waste, CO₂, microplastics",
            "Throwaway culture & overconsumption; poor quality",
            "Greenwashing – ‘looking’ sustainable without really changing",
        ],
        "notes": (
            "But the cons are serious. First, labour exploitation – poverty wages, child "
            "labour and unsafe factories, as the Earth.org article documents. Second, "
            "environmental harm – water pollution from dyeing, huge amounts of textile waste, "
            "CO2 emissions and microplastics. Third, it fuels a throwaway culture and "
            "overconsumption, and the clothes are often poor quality, so they do not last. "
            "And finally, many brands engage in greenwashing – pretending to be sustainable "
            "without really changing. So the cheap price hides a very high cost."
        ),
    },
    {
        "tag": "Discussion",
        "title": "How much responsibility does the consumer have?",
        "bullets": [
            "Shared responsibility – consumers create the demand…",
            "…but corporations & governments hold more power & information",
            "Fairness: not everyone can afford ethical alternatives",
            "Yet collective choices matter – we ‘vote with our wallets’",
        ],
        "notes": (
            "This leads to a key question: how much responsibility does the consumer have? "
            "Our view is that responsibility is shared. On one hand, consumers create the "
            "demand that drives the whole system. On the other hand, big corporations and "
            "governments hold far more structural power and far more information than the "
            "individual shopper. It is also a question of fairness: not everyone can afford "
            "ethical alternatives, so we should not put all the blame on the individual. "
            "Still, consumer choices matter – when enough people change what they buy, "
            "companies have to follow. We do 'vote with our wallets'."
        ),
    },
    {
        "tag": "Discussion",
        "title": "What can we as consumers do?",
        "bullets": [
            "Buy less, buy better – quality over quantity (cost-per-wear)",
            "Buy second-hand / swap / thrift; repair & care for clothes",
            "Support ethical brands; demand transparency",
            "Recycle & donate; back legislation → ‘slow fashion’",
        ],
        "notes": (
            "So what can we actually do as consumers? The simplest principle is: buy less, "
            "but buy better – choose quality over quantity and think about cost-per-wear. We "
            "can buy second-hand, swap clothes, or thrift, and we can repair and care for the "
            "clothes we already own to make them last. We can support ethical and sustainable "
            "brands and demand transparency about how clothes are made. And we can recycle, "
            "donate, and support legislation that protects workers and the environment. All "
            "of this is part of what is called 'slow fashion' – the opposite of the fast, "
            "throwaway model."
        ),
    },
    {
        "tag": "Wrap-up",
        "title": "Conclusion",
        "bullets": [
            "Fast fashion is a clear product of consumerism & materialism",
            "Earth.org uses logos + pathos on a base of institutional ethos",
            "It exposes the hidden human & environmental cost",
            "Responsibility is shared – but small changes, by many, add up",
        ],
        "notes": (
            "To conclude: fast fashion is a clear product of consumerism and materialism – "
            "our society's drive to keep buying, and our personal attachment to owning more. "
            "In its article, Earth.org uses a strong combination of logos and pathos, backed "
            "by solid institutional ethos, to expose the hidden human and environmental cost "
            "behind cheap clothing. And while responsibility is shared between consumers, "
            "companies and governments, we as consumers are not powerless – small changes, "
            "made by many people, really do add up. Thank you."
        ),
    },
    {
        "tag": "References",
        "title": "Sources",
        "bullets": [
            "“Fast Fashion: The Danger of Sweatshops”, Earth.org – earth.org/sweatshops/",
            "“The Greening of Throwaway Stuff”, BusinessLike, Systime",
            "(Add any further sources you used)",
        ],
        "notes": "Here are our sources.",
    },
    {
        "type": "closing",
        "title": "Thank you",
        "subtitle": "Any questions?",
        "notes": "Thank you for listening. We are happy to take any questions.",
    },
]


# ---------------------------------------------------------------------------
# PPTX helpers
# ---------------------------------------------------------------------------
def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if not line:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space_after=8, line_spacing=1.1):
    """runs: list of paragraphs; each paragraph is a list of (text, font, size, color, bold)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (text, font, size, color, bold) in para:
            r = p.add_run()
            r.text = text
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
    return tb


def add_content_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    # left accent bar
    add_rect(slide, Inches(0), Inches(0), Inches(0.16), Inches(7.5), PRIMARY)

    top = 0.55
    # eyebrow / tag
    if data.get("tag"):
        add_text(slide, Inches(0.7), Inches(top), Inches(11.8), Inches(0.4),
                 [[(data["tag"].upper(), BODY, 13, ACCENT, True)]])
        top += 0.42
    # title
    add_text(slide, Inches(0.7), Inches(top), Inches(12.0), Inches(1.0),
             [[(data["title"], HEAD, 30, PRIMARY, True)]])
    # underline rule
    add_rect(slide, Inches(0.72), Inches(top + 0.92), Inches(1.7), Inches(0.045), ACCENT)

    # body bullets
    paras = []
    for item in data.get("bullets", []):
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        if level == 0:
            paras.append([("▸  ", BODY, 16, ACCENT, True),
                          (text, BODY, 19, INK, False)])
        else:
            paras.append([("      –  ", BODY, 15, MUTED, False),
                          (text, BODY, 16, MUTED_DARK, False)])
    if paras:
        add_text(slide, Inches(0.95), Inches(top + 1.25), Inches(11.6), Inches(5.0),
                 paras, space_after=10, line_spacing=1.12)

    # footer
    add_text(slide, Inches(0.7), Inches(7.05), Inches(8.0), Inches(0.35),
             [[("Consumerism & Fast Fashion", BODY, 10, MUTED, False)]])
    idx = len(prs.slides._sldIdLst)  # current count
    add_text(slide, Inches(12.0), Inches(7.05), Inches(0.9), Inches(0.35),
             [[(str(idx), BODY, 10, MUTED, False)]], align=PP_ALIGN.RIGHT)

    slide.notes_slide.notes_text_frame.text = data.get("notes", "")
    return slide


def add_title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, PRIMARY)
    # accent block on the left
    add_rect(slide, Inches(0), Inches(0), Inches(0.35), Inches(7.5), ACCENT)
    # eyebrow
    add_text(slide, Inches(1.0), Inches(1.7), Inches(11.0), Inches(0.5),
             [[("GROUP PRESENTATION", BODY, 16, RGBColor(0xD8, 0xC9, 0xA8), True)]])
    # title
    add_text(slide, Inches(1.0), Inches(2.3), Inches(11.3), Inches(1.8),
             [[(data["title"], HEAD, 50, WHITE, True)]])
    # accent rule
    add_rect(slide, Inches(1.05), Inches(3.95), Inches(2.4), Inches(0.06), ACCENT)
    # subtitle
    add_text(slide, Inches(1.0), Inches(4.15), Inches(11.0), Inches(0.7),
             [[(data["subtitle"], HEAD, 24, RGBColor(0xCF, 0xDE, 0xD4), False)]])
    # meta
    meta_paras = [[(line, BODY, 15, RGBColor(0xB9, 0xCB, 0xC0), False)] for line in data["meta"]]
    add_text(slide, Inches(1.0), Inches(5.5), Inches(11.0), Inches(1.6),
             meta_paras, space_after=4, line_spacing=1.2)

    slide.notes_slide.notes_text_frame.text = data.get("notes", "")
    return slide


def add_closing_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, PRIMARY)
    add_rect(slide, Inches(0), Inches(0), Inches(0.35), Inches(7.5), ACCENT)
    add_text(slide, Inches(1.0), Inches(2.7), Inches(11.3), Inches(1.6),
             [[(data["title"], HEAD, 54, WHITE, True)]])
    add_rect(slide, Inches(1.05), Inches(4.35), Inches(2.4), Inches(0.06), ACCENT)
    add_text(slide, Inches(1.0), Inches(4.6), Inches(11.0), Inches(0.8),
             [[(data["subtitle"], HEAD, 26, RGBColor(0xCF, 0xDE, 0xD4), False)]])
    slide.notes_slide.notes_text_frame.text = data.get("notes", "")
    return slide


def add_pentagram_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)
    add_rect(slide, Inches(0), Inches(0), Inches(0.16), Inches(7.5), PRIMARY)

    top = 0.55
    add_text(slide, Inches(0.7), Inches(top), Inches(11.8), Inches(0.4),
             [[(data["tag"].upper(), BODY, 13, ACCENT, True)]])
    add_text(slide, Inches(0.7), Inches(top + 0.42), Inches(12.0), Inches(1.0),
             [[(data["title"], HEAD, 30, PRIMARY, True)]])
    add_rect(slide, Inches(0.72), Inches(top + 1.34), Inches(1.7), Inches(0.045), ACCENT)

    # short bullets on the left
    paras = [[("▸  ", BODY, 16, ACCENT, True), (b, BODY, 18, INK, False)]
             for b in data.get("bullets", [])]
    add_text(slide, Inches(0.95), Inches(2.55), Inches(4.7), Inches(4.0),
             paras, space_after=14, line_spacing=1.15)

    # ---- draw the pentagram on the right ----
    labels = [
        ("Topic", "Emne"),
        ("Receiver", "Modtager"),
        ("Language", "Sprog"),
        ("Circumstances", "Omstændigheder"),
        ("Sender", "Afsender"),
    ]
    cx, cy, r = 9.4, 4.25, 1.95  # inches
    pts = []
    for i in range(5):
        ang = math.radians(-90 + 72 * i)
        pts.append((cx + r * math.cos(ang), cy + r * math.sin(ang)))

    # connectors between every pair (pentagon + inner star)
    for i in range(5):
        for j in range(i + 1, 5):
            x1, y1 = pts[i]
            x2, y2 = pts[j]
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
            conn.line.color.rgb = SECONDARY
            conn.line.width = Pt(1.25)

    # center label
    add_text(slide, Inches(cx - 1.1), Inches(cy - 0.3), Inches(2.2), Inches(0.6),
             [[("the text", HEAD, 13, MUTED, False)],
              [("& its situation", HEAD, 13, MUTED, False)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.0)

    # nodes
    nw, nh = 2.0, 0.72
    for (en, da), (px, py) in zip(labels, pts):
        node = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(px - nw / 2), Inches(py - nh / 2), Inches(nw), Inches(nh))
        node.fill.solid()
        node.fill.fore_color.rgb = PRIMARY
        node.line.color.rgb = PRIMARY
        node.shadow.inherit = False
        tf = node.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 0.95
        r1 = p.add_run(); r1.text = en + "\n"
        r1.font.name = BODY; r1.font.size = Pt(15); r1.font.bold = True; r1.font.color.rgb = WHITE
        r2 = p.add_run(); r2.text = da
        r2.font.name = BODY; r2.font.size = Pt(11); r2.font.italic = True
        r2.font.color.rgb = RGBColor(0xC9, 0xD8, 0xCF)

    # footer
    add_text(slide, Inches(0.7), Inches(7.05), Inches(8.0), Inches(0.35),
             [[("Consumerism & Fast Fashion", BODY, 10, MUTED, False)]])
    idx = len(prs.slides._sldIdLst)
    add_text(slide, Inches(12.0), Inches(7.05), Inches(0.9), Inches(0.35),
             [[(str(idx), BODY, 10, MUTED, False)]], align=PP_ALIGN.RIGHT)

    slide.notes_slide.notes_text_frame.text = data.get("notes", "")
    return slide


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for data in SLIDES:
        t = data.get("type")
        if t == "title":
            add_title_slide(prs, data)
        elif t == "closing":
            add_closing_slide(prs, data)
        elif t == "pentagram":
            add_pentagram_slide(prs, data)
        else:
            add_content_slide(prs, data)
    out = os.path.join(HERE, "Fast_Fashion_Presentation.pptx")
    prs.save(out)
    return out


# ---------------------------------------------------------------------------
# DOCX speaker-notes manuscript
# ---------------------------------------------------------------------------
def _style_run(run, name=BODY, size=11, color=INK, bold=False, italic=False):
    run.font.name = name
    run.font.size = DPt(size)
    run.font.color.rgb = DRGB(color[0], color[1], color[2])
    run.font.bold = bold
    run.font.italic = italic


def build_docx():
    doc = Document()
    # base style
    normal = doc.styles["Normal"]
    normal.font.name = BODY
    normal.font.size = DPt(11)

    # ---- cover ----
    h = doc.add_paragraph()
    h.alignment = WD_ALIGN_PARAGRAPH.LEFT
    r = h.add_run("Consumerism & Fast Fashion")
    _style_run(r, name=HEAD, size=26, color=PRIMARY, bold=True)
    sub = doc.add_paragraph()
    _style_run(sub.add_run("Speaker notes / manuscript for the recorded presentation"),
               name=HEAD, size=14, color=MUTED_DARK, italic=True)

    intro = doc.add_paragraph()
    _style_run(intro.add_run(
        "How to use this document: each section below matches one slide. The text in "
        "“Manuscript” is written to be read aloud while you record – keep it natural, "
        "and feel free to adapt it to your own voice. The same notes are also in the "
        "Notes pane of the PowerPoint, so you can see them in Presenter view while you "
        "record."), size=11, color=INK)

    tip = doc.add_paragraph()
    _style_run(tip.add_run(
        "Tip (makkerskab): split the slides between you – e.g. one person takes the "
        "terms + pentagram, the other takes persuasion + discussion. Remember to insert "
        "your names on the title slide and to double-check the quotes/figures against "
        "your own copy of the two texts."), size=11, color=MUTED_DARK, italic=True)

    rec = doc.add_paragraph()
    _style_run(rec.add_run(
        "Recording in PowerPoint: Slide Show → Record → From Beginning. Speak over each "
        "slide, then File → Export → Create a Video (or save and the recording stays "
        "embedded)."), size=11, color=MUTED_DARK, italic=True)

    doc.add_paragraph()  # spacer

    # ---- per-slide notes ----
    for i, data in enumerate(SLIDES, start=1):
        title = data.get("title", "")
        head = doc.add_paragraph()
        _style_run(head.add_run(f"Slide {i} – {title}"),
                   name=HEAD, size=15, color=PRIMARY, bold=True)
        # thin spacing line
        if data.get("tag"):
            tagp = doc.add_paragraph()
            _style_run(tagp.add_run(data["tag"].upper()), size=9, color=ACCENT, bold=True)

        # on-slide bullets summary
        bullets = data.get("bullets")
        if bullets:
            lbl = doc.add_paragraph()
            _style_run(lbl.add_run("On the slide:"), size=10, color=MUTED, bold=True)
            for item in bullets:
                if isinstance(item, tuple):
                    level, text = item
                else:
                    level, text = 0, item
                bp = doc.add_paragraph(style="List Bullet" if level == 0 else "List Bullet 2")
                _style_run(bp.add_run(text), size=10.5, color=MUTED_DARK)

        # manuscript
        mlbl = doc.add_paragraph()
        _style_run(mlbl.add_run("Manuscript:"), size=10, color=MUTED, bold=True)
        mp = doc.add_paragraph()
        _style_run(mp.add_run(data.get("notes", "")), size=11.5, color=INK)

        doc.add_paragraph()  # spacer between slides

    out = os.path.join(HERE, "Fast_Fashion_Speaker_Notes.docx")
    doc.save(out)
    return out


if __name__ == "__main__":
    p = build_pptx()
    d = build_docx()
    print("Wrote:", p)
    print("Wrote:", d)
